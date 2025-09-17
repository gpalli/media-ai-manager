#!/usr/bin/env python3
"""
Media file processor for MediaMind AI.
Handles extraction of metadata and content from various media file types.
Based on the original implementation but extended for images and videos.
"""

import os
import hashlib
import mimetypes
import logging
from pathlib import Path
from datetime import datetime
from typing import Dict, List, Optional, Any, Tuple
from dataclasses import dataclass
import json

# Image processing
from PIL import Image, ExifTags
from PIL.ExifTags import TAGS
import exifread

# Video processing
import cv2
import numpy as np

# Optional video processing
try:
    from moviepy.editor import VideoFileClip
    MOVIEPY_AVAILABLE = True
except ImportError:
    MOVIEPY_AVAILABLE = False
    VideoFileClip = None

# Document processing (from original implementation)
try:
    import PyPDF2
    PYPDF2_AVAILABLE = True
except ImportError:
    PYPDF2_AVAILABLE = False
    PyPDF2 = None

try:
    import docx
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    docx = None

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    Presentation = None

# Utilities
try:
    import magic
    MAGIC_AVAILABLE = True
except ImportError:
    MAGIC_AVAILABLE = False
    magic = None

from tqdm import tqdm

logger = logging.getLogger(__name__)

@dataclass
class MediaFile:
    """Represents a media file with all its metadata."""
    file_path: str
    file_name: str
    file_size: int
    file_hash: str
    mime_type: str
    file_type: str  # 'image', 'video', 'document'
    created_date: datetime
    modified_date: datetime
    
    # Media-specific metadata
    width: Optional[int] = None
    height: Optional[int] = None
    duration: Optional[float] = None
    fps: Optional[float] = None
    bitrate: Optional[int] = None
    
    # EXIF and technical metadata
    exif_data: Optional[Dict] = None
    camera_info: Optional[Dict] = None
    gps_data: Optional[Dict] = None
    
    # Content data
    text_content: Optional[str] = None
    thumbnail_path: Optional[str] = None
    
    # AI-generated metadata (will be filled by AI analyzer)
    ai_description: Optional[str] = None
    ai_tags: List[str] = None
    detected_objects: List[str] = None
    scene_type: Optional[str] = None
    extracted_text: Optional[str] = None
    
    def __post_init__(self):
        if self.ai_tags is None:
            self.ai_tags = []

class MediaProcessor:
    """Processes various media file types and extracts metadata."""
    
    def __init__(self, config: Dict):
        self.config = config
        self.media_config = config.get('media', {})
        self.supported_formats = self.media_config.get('supported_formats', {})
        
        # Create thumbnail directory
        self.thumbnail_dir = Path("cache/thumbnails")
        self.thumbnail_dir.mkdir(parents=True, exist_ok=True)
    
    def process_file(self, file_path: str) -> Optional[MediaFile]:
        """Process a single media file and extract all metadata."""
        try:
            if not os.path.exists(file_path):
                logger.error(f"File does not exist: {file_path}")
                return None
            
            # Check file size
            file_size = os.path.getsize(file_path)
            max_size = self.media_config.get('max_file_size_mb', 500) * 1024 * 1024
            if file_size > max_size:
                logger.warning(f"File too large ({file_size / 1024 / 1024:.1f}MB): {file_path}")
                return None
            
            # Get basic file info
            stat = os.stat(file_path)
            file_hash = self._calculate_file_hash(file_path)
            
            if MAGIC_AVAILABLE:
                mime_type = magic.from_file(file_path, mime=True)
            else:
                # Fallback to mimetypes module
                import mimetypes
                mime_type, _ = mimetypes.guess_type(file_path)
                if not mime_type:
                    mime_type = 'application/octet-stream'
            
            file_type = self._determine_file_type(file_path, mime_type)
            
            # Create MediaFile object
            media_file = MediaFile(
                file_path=file_path,
                file_name=os.path.basename(file_path),
                file_size=file_size,
                file_hash=file_hash,
                mime_type=mime_type,
                file_type=file_type,
                created_date=datetime.fromtimestamp(stat.st_ctime),
                modified_date=datetime.fromtimestamp(stat.st_mtime)
            )
            
            # Process based on file type
            if file_type == 'image':
                media_file = self._process_image(media_file)
            elif file_type == 'video':
                media_file = self._process_video(media_file)
            elif file_type == 'document':
                media_file = self._process_document(media_file)
            else:
                logger.warning(f"Unsupported file type: {file_path}")
                return None
            
            # Generate thumbnail
            media_file.thumbnail_path = self._generate_thumbnail(media_file)
            
            return media_file
            
        except Exception as e:
            logger.error(f"Error processing file {file_path}: {e}")
            return None
    
    def _determine_file_type(self, file_path: str, mime_type: str) -> str:
        """Determine the type of media file."""
        file_ext = Path(file_path).suffix.lower()
        
        if file_ext in self.supported_formats.get('images', []):
            return 'image'
        elif file_ext in self.supported_formats.get('videos', []):
            return 'video'
        elif file_ext in self.supported_formats.get('documents', []):
            return 'document'
        else:
            return 'unknown'
    
    def _calculate_file_hash(self, file_path: str) -> str:
        """Calculate SHA-256 hash of file for deduplication."""
        hash_sha256 = hashlib.sha256()
        try:
            with open(file_path, "rb") as f:
                for chunk in iter(lambda: f.read(4096), b""):
                    hash_sha256.update(chunk)
            return hash_sha256.hexdigest()
        except Exception as e:
            logger.error(f"Error calculating hash for {file_path}: {e}")
            return ""
    
    def _process_image(self, media_file: MediaFile) -> MediaFile:
        """Process image files and extract metadata."""
        try:
            with Image.open(media_file.file_path) as img:
                # Get dimensions
                media_file.width, media_file.height = img.size
                
                # Extract EXIF data
                media_file.exif_data = self._extract_exif_data(img)
                
                # Extract camera info
                media_file.camera_info = self._extract_camera_info(media_file.exif_data)
                
                # Extract GPS data
                media_file.gps_data = self._extract_gps_data(media_file.exif_data)
                
                logger.debug(f"Processed image: {media_file.file_name} ({media_file.width}x{media_file.height})")
                
        except Exception as e:
            logger.error(f"Error processing image {media_file.file_path}: {e}")
        
        return media_file
    
    def _process_video(self, media_file: MediaFile) -> MediaFile:
        """Process video files and extract metadata."""
        try:
            # Use OpenCV for basic video info
            cap = cv2.VideoCapture(media_file.file_path)
            if cap.isOpened():
                media_file.width = int(cap.get(cv2.CAP_PROP_FRAME_WIDTH))
                media_file.height = int(cap.get(cv2.CAP_PROP_FRAME_HEIGHT))
                media_file.fps = cap.get(cv2.CAP_PROP_FPS)
                frame_count = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
                if media_file.fps > 0:
                    media_file.duration = frame_count / media_file.fps
                cap.release()
            
            # Use MoviePy for additional metadata (if available)
            if MOVIEPY_AVAILABLE:
                try:
                    with VideoFileClip(media_file.file_path) as clip:
                        media_file.duration = clip.duration
                        media_file.fps = clip.fps
                        # Get bitrate if available
                        if hasattr(clip, 'bitrate'):
                            media_file.bitrate = clip.bitrate
                except Exception as e:
                    logger.debug(f"MoviePy processing failed for {media_file.file_path}: {e}")
            else:
                logger.debug("MoviePy not available, using OpenCV only for video processing")
            
            logger.debug(f"Processed video: {media_file.file_name} ({media_file.width}x{media_file.height}, {media_file.duration:.1f}s)")
            
        except Exception as e:
            logger.error(f"Error processing video {media_file.file_path}: {e}")
        
        return media_file
    
    def _process_document(self, media_file: MediaFile) -> MediaFile:
        """Process document files and extract text content."""
        try:
            file_ext = Path(media_file.file_path).suffix.lower()
            
            if file_ext == '.pdf':
                media_file.text_content = self._extract_pdf_text(media_file.file_path)
            elif file_ext == '.docx':
                media_file.text_content = self._extract_docx_text(media_file.file_path)
            elif file_ext == '.pptx':
                media_file.text_content = self._extract_pptx_text(media_file.file_path)
            elif file_ext in ['.txt', '.md']:
                media_file.text_content = self._extract_text_file(media_file.file_path)
            
            logger.debug(f"Processed document: {media_file.file_name} ({len(media_file.text_content or '')} chars)")
            
        except Exception as e:
            logger.error(f"Error processing document {media_file.file_path}: {e}")
        
        return media_file
    
    def _extract_exif_data(self, img: Image.Image) -> Dict:
        """Extract EXIF data from image."""
        exif_data = {}
        try:
            if hasattr(img, '_getexif') and img._getexif() is not None:
                exif = img._getexif()
                for tag_id, value in exif.items():
                    tag = TAGS.get(tag_id, tag_id)
                    exif_data[tag] = value
        except Exception as e:
            logger.debug(f"Error extracting EXIF data: {e}")
        
        return exif_data
    
    def _extract_camera_info(self, exif_data: Dict) -> Dict:
        """Extract camera information from EXIF data."""
        camera_info = {}
        
        if not exif_data:
            return camera_info
        
        # Camera make and model
        camera_info['make'] = exif_data.get('Make', '')
        camera_info['model'] = exif_data.get('Model', '')
        
        # Camera settings
        camera_info['iso'] = exif_data.get('ISOSpeedRatings', '')
        camera_info['focal_length'] = exif_data.get('FocalLength', '')
        camera_info['aperture'] = exif_data.get('FNumber', '')
        camera_info['shutter_speed'] = exif_data.get('ExposureTime', '')
        
        # Flash and other settings
        camera_info['flash'] = exif_data.get('Flash', '')
        camera_info['white_balance'] = exif_data.get('WhiteBalance', '')
        
        return camera_info
    
    def _extract_gps_data(self, exif_data: Dict) -> Dict:
        """Extract GPS data from EXIF data."""
        gps_data = {}
        
        if not exif_data:
            return gps_data
        
        # GPS coordinates
        gps_latitude = exif_data.get('GPSInfo', {}).get('GPSLatitude')
        gps_latitude_ref = exif_data.get('GPSInfo', {}).get('GPSLatitudeRef')
        gps_longitude = exif_data.get('GPSInfo', {}).get('GPSLongitude')
        gps_longitude_ref = exif_data.get('GPSInfo', {}).get('GPSLongitudeRef')
        
        if gps_latitude and gps_longitude:
            # Convert to decimal degrees
            lat = self._convert_to_degrees(gps_latitude)
            lon = self._convert_to_degrees(gps_longitude)
            
            if gps_latitude_ref == 'S':
                lat = -lat
            if gps_longitude_ref == 'W':
                lon = -lon
            
            gps_data['latitude'] = lat
            gps_data['longitude'] = lon
        
        return gps_data
    
    def _convert_to_degrees(self, value):
        """Convert GPS coordinates to decimal degrees."""
        d, m, s = value
        return d + (m / 60.0) + (s / 3600.0)
    
    def _extract_pdf_text(self, file_path: str) -> str:
        """Extract text from PDF file."""
        if not PYPDF2_AVAILABLE:
            logger.warning("PyPDF2 not available, cannot extract PDF text")
            return ""
        
        try:
            with open(file_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                text = ' '.join([page.extract_text() for page in reader.pages])
                return text
        except Exception as e:
            logger.error(f"Error extracting PDF text: {e}")
            return ""
    
    def _extract_docx_text(self, file_path: str) -> str:
        """Extract text from DOCX file."""
        if not DOCX_AVAILABLE:
            logger.warning("python-docx not available, cannot extract DOCX text")
            return ""
        
        try:
            doc = docx.Document(file_path)
            text = ' '.join([para.text for para in doc.paragraphs])
            return text
        except Exception as e:
            logger.error(f"Error extracting DOCX text: {e}")
            return ""
    
    def _extract_pptx_text(self, file_path: str) -> str:
        """Extract text from PPTX file."""
        if not PPTX_AVAILABLE:
            logger.warning("python-pptx not available, cannot extract PPTX text")
            return ""
        
        try:
            prs = Presentation(file_path)
            text = ' '.join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, 'text')])
            return text
        except Exception as e:
            logger.error(f"Error extracting PPTX text: {e}")
            return ""
    
    def _extract_text_file(self, file_path: str) -> str:
        """Extract text from plain text file."""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except Exception as e:
            logger.error(f"Error reading text file: {e}")
            return ""
    
    def _generate_thumbnail(self, media_file: MediaFile) -> Optional[str]:
        """Generate thumbnail for media file."""
        try:
            thumbnail_name = f"{media_file.file_hash}.jpg"
            thumbnail_path = self.thumbnail_dir / thumbnail_name
            
            if thumbnail_path.exists():
                return str(thumbnail_path)
            
            if media_file.file_type == 'image':
                return self._generate_image_thumbnail(media_file, thumbnail_path)
            elif media_file.file_type == 'video':
                return self._generate_video_thumbnail(media_file, thumbnail_path)
            
        except Exception as e:
            logger.error(f"Error generating thumbnail for {media_file.file_path}: {e}")
        
        return None
    
    def _generate_image_thumbnail(self, media_file: MediaFile, thumbnail_path: Path) -> str:
        """Generate thumbnail for image file."""
        try:
            with Image.open(media_file.file_path) as img:
                # Convert to RGB if necessary
                if img.mode != 'RGB':
                    img = img.convert('RGB')
                
                # Create thumbnail
                img.thumbnail(self.media_config.get('thumbnail_size', [300, 300]))
                img.save(thumbnail_path, 'JPEG', quality=85)
                
            return str(thumbnail_path)
        except Exception as e:
            logger.error(f"Error generating image thumbnail: {e}")
            return None
    
    def _generate_video_thumbnail(self, media_file: MediaFile, thumbnail_path: Path) -> str:
        """Generate thumbnail for video file."""
        try:
            cap = cv2.VideoCapture(media_file.file_path)
            if cap.isOpened():
                # Get frame from middle of video
                frame_count = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
                middle_frame = frame_count // 2
                cap.set(cv2.CAP_PROP_POS_FRAMES, middle_frame)
                
                ret, frame = cap.read()
                if ret:
                    # Convert BGR to RGB
                    frame_rgb = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
                    
                    # Resize to thumbnail size
                    thumbnail_size = self.media_config.get('thumbnail_size', [300, 300])
                    frame_resized = cv2.resize(frame_rgb, tuple(thumbnail_size))
                    
                    # Save as JPEG
                    img = Image.fromarray(frame_resized)
                    img.save(thumbnail_path, 'JPEG', quality=85)
                
                cap.release()
                return str(thumbnail_path)
        except Exception as e:
            logger.error(f"Error generating video thumbnail: {e}")
        
        return None
    
    def discover_media_files(self, scan_paths: List[str]) -> List[str]:
        """Discover all media files in the specified paths."""
        media_files = []
        excluded_paths = self.media_config.get('excluded_paths', [])
        
        all_formats = []
        for formats in self.supported_formats.values():
            all_formats.extend(formats)
        
        for scan_path in scan_paths:
            if not os.path.exists(scan_path):
                logger.warning(f"Scan path does not exist: {scan_path}")
                continue
            
            logger.info(f"Scanning directory: {scan_path}")
            
            for root, dirs, files in os.walk(scan_path):
                # Skip excluded directories
                if any(excluded in root for excluded in excluded_paths):
                    continue
                
                # Filter directories to avoid scanning
                dirs[:] = [d for d in dirs if not any(excluded in os.path.join(root, d) for excluded in excluded_paths)]
                
                for file in files:
                    file_path = os.path.join(root, file)
                    file_ext = Path(file).suffix.lower()
                    
                    if file_ext in all_formats:
                        try:
                            file_size = os.path.getsize(file_path)
                            max_size = self.media_config.get('max_file_size_mb', 500) * 1024 * 1024
                            if file_size <= max_size:
                                media_files.append(file_path)
                        except (OSError, IOError):
                            continue
        
        logger.info(f"Discovered {len(media_files)} media files")
        return media_files
